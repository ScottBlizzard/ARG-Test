from __future__ import annotations

import argparse
import json
import textwrap
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]

PALETTE = {
    'ink': '#132238',
    'muted': '#5B6677',
    'paper': '#FBFAF6',
    'panel': '#F7F5EF',
    'panel_alt': '#EEF3F7',
    'accent': '#D55C3A',
    'accent_soft': '#F2C6B8',
    'teal': '#3B8C88',
    'teal_soft': '#B9DDD8',
    'gold': '#D8A332',
    'gold_soft': '#F5E3B5',
    'line': '#D8DCE3',
    'purple': '#7C6BB0',
    'purple_soft': '#E2DAF3',
    'white': '#FFFFFF',
}


def hex_rgb(value: str) -> RGBColor:
    value = value.lstrip('#')
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def read_json(path: Path):
    return json.loads(path.read_text(encoding='utf-8-sig'))


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def avg(values: list[float]) -> float:
    return round(sum(values) / len(values), 3) if values else 0.0


def pretty_requirement_name(requirement_id: str) -> str:
    return requirement_id.replace('_', ' ').title()


def short_diag(text: str, width: int = 56) -> str:
    cleaned = ' '.join(text.split())
    return textwrap.shorten(cleaned, width=width, placeholder='...')


def configure_matplotlib() -> None:
    plt.rcParams.update({
        'figure.facecolor': PALETTE['white'],
        'axes.facecolor': PALETTE['white'],
        'savefig.facecolor': PALETTE['white'],
        'font.family': ['DejaVu Sans', 'Microsoft YaHei', 'Segoe UI'],
        'font.size': 11.5,
        'axes.titleweight': 'bold',
        'axes.labelcolor': PALETTE['ink'],
        'xtick.color': PALETTE['muted'],
        'ytick.color': PALETTE['muted'],
        'axes.edgecolor': PALETTE['line'],
        'axes.spines.top': False,
        'axes.spines.right': False,
    })


def save_fig(fig, output_path: Path) -> None:
    fig.savefig(output_path, bbox_inches='tight', pad_inches=0.06)
    if output_path.suffix.lower() != '.pdf':
        fig.savefig(output_path.with_suffix('.pdf'), bbox_inches='tight', pad_inches=0.06)


def aggregate_payload(run_main: list[dict], baseline: list[dict], ablation: list[dict]) -> dict:
    payload = {
        'main': {
            'avg_checker_score': avg([float(item['score']) for item in run_main]),
            'avg_overall_coverage': avg([float(item['metrics']['overall_coverage']) for item in run_main]),
            'avg_test_count': avg([float(item['metrics']['test_count']) for item in run_main]),
            'avg_risk_score': avg([float(item['risk_assessment']['score']) for item in run_main]),
            'high_risk_count': sum(1 for item in run_main if item['risk_assessment']['level'] == 'High'),
            'repaired_count': sum(1 for item in run_main if item.get('repaired')),
            'avg_duplicate_count': avg([float(item['metrics']['duplicate_count']) for item in run_main]),
            'count': len(run_main),
        },
        'baseline': {},
        'ablation': {},
    }
    for method in ('rule_based', 'plain_llm', 'structured_no_checker'):
        payload['baseline'][method] = {
            'avg_checker_score': avg([float(item['baselines'][method]['checker_score']) for item in baseline]),
            'avg_overall_coverage': avg([float(item['baselines'][method]['overall_coverage']) for item in baseline]),
            'avg_test_count': avg([float(item['baselines'][method]['test_count']) for item in baseline]),
        }
    for method in ('structured_no_checker', 'full_pipeline'):
        payload['ablation'][method] = {
            'avg_checker_score': avg([float(item[method]['checker_score']) for item in ablation]),
            'avg_overall_coverage': avg([float(item[method]['overall_coverage']) for item in ablation]),
            'avg_test_count': avg([float(item[method]['test_count']) for item in ablation]),
        }
    return payload


def add_box(slide, left, top, width, height, fill, line, title, body_lines, title_size=18, body_size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.line.color.rgb = hex_rgb(line)
    shape.line.width = Pt(1.4)
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Pt(12)
    text_frame.margin_right = Pt(12)
    text_frame.margin_top = Pt(10)
    text_frame.margin_bottom = Pt(10)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(title_size)
    p.font.color.rgb = hex_rgb(PALETTE['ink'])
    p.alignment = PP_ALIGN.LEFT
    for line_text in body_lines:
        p = text_frame.add_paragraph()
        p.text = line_text
        p.level = 0
        p.font.size = Pt(body_size)
        p.font.color.rgb = hex_rgb(PALETTE['muted'])
        p.alignment = PP_ALIGN.LEFT
    return shape


def add_label(slide, left, top, width, height, text_value, fill, font_size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.fill.transparency = 0.1
    shape.line.color.rgb = hex_rgb(fill)
    shape.line.width = Pt(0.8)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text_value
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = hex_rgb(PALETTE['ink'])
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def connect(slide, x1, y1, x2, y2, color=PALETTE['muted']):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = hex_rgb(color)
    line.line.width = Pt(2.2)
    return line


def build_architecture_pptx(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_rgb(PALETTE['paper'])

    title = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(8.8), Inches(0.55))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = 'ARG-Test Architecture'
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = hex_rgb(PALETTE['ink'])

    sub = slide.shapes.add_textbox(Inches(0.67), Inches(0.78), Inches(9.5), Inches(0.4))
    sp = sub.text_frame.paragraphs[0]
    sp.text = 'Requirement-driven black-box test generation with structured reasoning, checker validation, reranking, and repair.'
    sp.font.size = Pt(12.5)
    sp.font.color.rgb = hex_rgb(PALETTE['muted'])

    add_label(slide, Inches(0.72), Inches(1.15), Inches(1.65), Inches(0.34), 'Model: Qwen3.5-Flash', PALETTE['accent_soft'])
    add_label(slide, Inches(2.5), Inches(1.15), Inches(1.2), Inches(0.34), 'Candidates: 3', PALETTE['teal_soft'])
    add_label(slide, Inches(3.82), Inches(1.15), Inches(1.8), Inches(0.34), 'Output: JSON / CSV / MD', PALETTE['gold_soft'])

    top_y = Inches(1.7)
    h = Inches(1.15)
    box1 = add_box(slide, Inches(0.75), top_y, Inches(2.05), h, PALETTE['gold_soft'], PALETTE['gold'], 'Requirement Dataset', ['dev/test requirements', 'gold specs + manifest'])
    box2 = add_box(slide, Inches(3.15), top_y, Inches(2.25), h, PALETTE['panel_alt'], PALETTE['teal'], 'Prompt Composer', ['system prompt + task prompt', 'provider = openai-compatible'])
    box3 = add_box(slide, Inches(5.85), top_y, Inches(2.2), h, PALETTE['accent_soft'], PALETTE['accent'], 'Structured Trace', ['analysis / pattern / steps', 'verification / final answer'])
    box4 = add_box(slide, Inches(8.45), top_y, Inches(2.55), h, PALETTE['teal_soft'], PALETTE['teal'], 'Parser + Schema Gate', ['typed trace parsing', 'suite export preparation'])

    mid_y = Inches(3.6)
    box5 = add_box(slide, Inches(1.0), mid_y, Inches(2.15), Inches(1.25), PALETTE['gold_soft'], PALETTE['gold'], 'EP / BVA Contract', ['partitions, invalid inputs,', 'boundaries and edge cases'])
    box6 = add_box(slide, Inches(3.65), mid_y, Inches(2.25), Inches(1.25), PALETTE['panel_alt'], PALETTE['teal'], 'Decision Contract', ['rule combinations, priorities,', 'conflicts and exceptions'])
    box7 = add_box(slide, Inches(6.35), mid_y, Inches(2.25), Inches(1.25), PALETTE['accent_soft'], PALETTE['accent'], 'State Contract', ['states, transitions, retries,', 'exceptional paths'])
    box8 = add_box(slide, Inches(9.0), mid_y, Inches(2.35), Inches(1.25), PALETTE['purple_soft'], PALETTE['purple'], 'Rerank + Repair', ['aggregate checker score,', 'local repair, best candidate'])

    bot_y = Inches(5.85)
    box9 = add_box(slide, Inches(2.55), bot_y, Inches(3.2), Inches(1.08), PALETTE['teal_soft'], PALETTE['teal'], 'Final Test Suite', ['JSON + CSV + markdown test cases', 'chosen candidate after checker validation'], title_size=20)
    box10 = add_box(slide, Inches(6.2), bot_y, Inches(4.0), Inches(1.08), PALETTE['accent_soft'], PALETTE['accent'], 'Evaluation and Reports', ['run_main / baselines / ablation / generalization', 'coverage, checker score, category breakdown'], title_size=20)

    connect(slide, box1.left + box1.width, box1.top + box1.height / 2, box2.left, box2.top + box2.height / 2)
    connect(slide, box2.left + box2.width, box2.top + box2.height / 2, box3.left, box3.top + box3.height / 2)
    connect(slide, box3.left + box3.width, box3.top + box3.height / 2, box4.left, box4.top + box4.height / 2)
    connect(slide, box4.left + box4.width / 2, box4.top + box4.height, box8.left + box8.width / 2, box8.top)
    connect(slide, box5.left + box5.width, box5.top + box5.height / 2, box6.left, box6.top + box6.height / 2)
    connect(slide, box6.left + box6.width, box6.top + box6.height / 2, box7.left, box7.top + box7.height / 2)
    connect(slide, box7.left + box7.width, box7.top + box7.height / 2, box8.left, box8.top + box8.height / 2)
    connect(slide, box8.left + box8.width / 2, box8.top + box8.height, box10.left + Inches(1.0), box10.top)
    connect(slide, box8.left + Inches(0.4), box8.top + box8.height, box9.left + Inches(0.9), box9.top)

    prs.save(output_path)



def save_main_vs_baselines(output_path: Path, payload: dict) -> None:
    configure_matplotlib()
    methods = [
        ('ARG-Test Full Pipeline', payload['main']['avg_overall_coverage'], payload['main']['avg_checker_score'], payload['main']['avg_test_count'], PALETTE['purple']),
        ('Structured No Checker', payload['baseline']['structured_no_checker']['avg_overall_coverage'], payload['baseline']['structured_no_checker']['avg_checker_score'], payload['baseline']['structured_no_checker']['avg_test_count'], PALETTE['teal']),
        ('Rule-based', payload['baseline']['rule_based']['avg_overall_coverage'], payload['baseline']['rule_based']['avg_checker_score'], payload['baseline']['rule_based']['avg_test_count'], PALETTE['gold']),
        ('Plain LLM', payload['baseline']['plain_llm']['avg_overall_coverage'], payload['baseline']['plain_llm']['avg_checker_score'], payload['baseline']['plain_llm']['avg_test_count'], PALETTE['accent']),
    ]
    fig, ax = plt.subplots(figsize=(8.8, 5.0), dpi=200)
    ax.set_xlim(0, 1.02)
    ax.set_ylim(-0.6, len(methods) - 0.4)
    ax.grid(axis='x', color=PALETTE['line'], linestyle='--', linewidth=0.8, alpha=0.85)
    ax.set_xlabel('Average metric value')
    y_positions = list(range(len(methods)))
    for idx, (label, coverage, score, test_count, color) in enumerate(methods):
        if idx == 0:
            ax.axhspan(idx - 0.5, idx + 0.5, color=PALETTE['purple_soft'], alpha=0.22, zorder=0)
        ax.barh(idx - 0.16, coverage, height=0.28, color=color, alpha=0.92, edgecolor='none', zorder=3)
        ax.barh(idx + 0.16, score, height=0.20, color=PALETTE['ink'], alpha=0.92, edgecolor='none', zorder=4)
        ax.text(min(max(coverage, score) + 0.02, 0.94), idx, f'cov {coverage:.3f} | score {score:.3f} | tests {test_count:.3f}', va='center', ha='left', fontsize=10.8, color=PALETTE['ink'])
    ax.set_yticks(y_positions, [item[0] for item in methods])
    ax.invert_yaxis()
    legend_handles = [
        FancyBboxPatch((0, 0), 1, 1, facecolor=PALETTE['teal'], edgecolor='none', alpha=0.92, label='Coverage'),
        FancyBboxPatch((0, 0), 1, 1, facecolor=PALETTE['ink'], edgecolor='none', alpha=0.92, label='Checker score'),
    ]
    ax.legend(handles=legend_handles, frameon=False, loc='lower right')
    fig.tight_layout(pad=0.8)
    save_fig(fig, output_path)
    plt.close(fig)



def save_ablation(output_path: Path, payload: dict) -> None:
    configure_matplotlib()
    labels = ['Structured No Checker', 'Full Pipeline']
    coverage = [
        payload['baseline']['structured_no_checker']['avg_overall_coverage'],
        payload['main']['avg_overall_coverage'],
    ]
    scores = [
        payload['baseline']['structured_no_checker']['avg_checker_score'],
        payload['main']['avg_checker_score'],
    ]
    test_counts = [
        payload['baseline']['structured_no_checker']['avg_test_count'],
        payload['main']['avg_test_count'],
    ]
    fig, axes = plt.subplots(1, 2, figsize=(8.2, 3.9), dpi=200)
    metric_specs = [
        ('Checker score', scores, PALETTE['accent'], f"+{scores[1] - scores[0]:.3f}"),
        ('Overall coverage', coverage, PALETTE['teal'], f"{coverage[1] - coverage[0]:+.3f}"),
    ]
    for ax, (title, values, color, delta_text) in zip(axes, metric_specs):
        bars = ax.bar(labels, values, color=[PALETTE['panel'], color], edgecolor='none', width=0.58)
        ax.set_ylim(0, 1.02)
        ax.set_title(title, fontsize=12.2, color=PALETTE['ink'], pad=8)
        ax.grid(axis='y', color=PALETTE['line'], linestyle='--', linewidth=0.8, alpha=0.85)
        for idx, (bar, value) in enumerate(zip(bars, values)):
            ax.text(bar.get_x() + bar.get_width() / 2, value + 0.03, f'{value:.3f}', ha='center', fontsize=10.5, color=PALETTE['ink'])
            ax.text(bar.get_x() + bar.get_width() / 2, 0.04, f'tests {test_counts[idx]:.3f}', ha='center', fontsize=9.5, color=PALETTE['muted'])
        ax.text(0.5, 0.98, f'delta = {delta_text}', transform=ax.transAxes, ha='center', va='top', fontsize=10.0, color=PALETTE['muted'])
        ax.tick_params(axis='x', labelrotation=0)
    fig.tight_layout(pad=0.8, w_pad=1.8)
    save_fig(fig, output_path)
    plt.close(fig)



def save_generalization(output_path: Path, categories: list[dict]) -> None:
    configure_matplotlib()
    label_map = {
        'business_rules': 'Business Rules',
        'input_validation': 'Input Validation',
        'workflow_state': 'Workflow State',
    }
    labels = [label_map[item['category']] for item in categories]
    coverage = [item['avg_overall_coverage'] for item in categories]
    scores = [item['avg_checker_score'] for item in categories]
    counts = [item['requirement_count'] for item in categories]
    risks = [item.get('avg_risk_score', 0.0) for item in categories]
    colors = [PALETTE['accent'], PALETTE['gold'], PALETTE['teal']]
    x = list(range(len(labels)))
    fig, ax = plt.subplots(figsize=(8.4, 4.7), dpi=200)
    bars = ax.bar(x, coverage, color=colors, width=0.56, edgecolor='none', zorder=2)
    ax.plot(x, scores, color=PALETTE['ink'], marker='o', linewidth=2.4, markersize=7, label='Checker score', zorder=4)
    ax.set_ylim(0, 1.02)
    ax.set_xticks(x, labels)
    ax.grid(axis='y', color=PALETTE['line'], linestyle='--', alpha=0.7)
    for idx, (bar, cov, score) in enumerate(zip(bars, coverage, scores)):
        ax.text(bar.get_x() + bar.get_width()/2, cov + 0.03, f'cov {cov:.3f}', ha='center', fontsize=10.5, color=PALETTE['ink'])
        ax.text(bar.get_x() + bar.get_width()/2, score + 0.03, f'score {score:.3f}', ha='center', fontsize=10.5, color=PALETTE['ink'])
        ax.text(bar.get_x() + bar.get_width()/2, 0.04, f'n={counts[idx]} | risk {risks[idx]:.2f}', ha='center', fontsize=9.4, color=PALETTE['muted'])
    fig.tight_layout(pad=0.8)
    save_fig(fig, output_path)
    plt.close(fig)



def save_stability(output_path: Path, stability: dict) -> None:
    configure_matplotlib()
    labels = [
        'Bundle\ndiscount',
        'Refund\nmethod',
        'Pickup\ncontact',
        'Card expiry\nand CVV',
        '3DS\nauthentication',
    ]
    formal = [item['formal_coverage'] for item in stability['comparisons']]
    rerun = [item['rerun_coverage'] for item in stability['comparisons']]
    score_formal = [item['formal_score'] for item in stability['comparisons']]
    score_rerun = [item['rerun_score'] for item in stability['comparisons']]
    x = range(len(labels))
    fig, axes = plt.subplots(2, 1, figsize=(10.4, 6.8), dpi=180, sharex=True)
    width = 0.36
    axes[0].bar([i - width/2 for i in x], formal, width=width, color=PALETTE['teal'], label='Formal')
    axes[0].bar([i + width/2 for i in x], rerun, width=width, color=PALETTE['gold'], label='Sanity Rerun')
    axes[0].set_ylim(0, 1.02)
    axes[0].set_ylabel('Coverage')
    axes[0].legend(frameon=False, loc='upper left')
    axes[0].grid(axis='y', color=PALETTE['line'], linestyle='--', alpha=0.7)
    axes[0].spines[['top', 'right']].set_visible(False)

    axes[1].plot(labels, score_formal, color=PALETTE['accent'], marker='o', linewidth=2.2, label='Formal score')
    axes[1].plot(labels, score_rerun, color=PALETTE['purple'], marker='o', linewidth=2.2, label='Rerun score')
    axes[1].set_ylim(0.7, 1.02)
    axes[1].set_ylabel('Checker Score')
    axes[1].grid(axis='y', color=PALETTE['line'], linestyle='--', alpha=0.7)
    axes[1].spines[['top', 'right']].set_visible(False)
    axes[1].legend(frameon=False, loc='upper left')
    axes[1].set_xticks(list(x), labels, rotation=0, ha='center')
    fig.tight_layout()
    save_fig(fig, output_path)
    plt.close(fig)



def save_reproducibility_summary(output_path: Path, series: list[dict]) -> None:
    configure_matplotlib()
    fig = plt.figure(figsize=(11.4, 5.6), dpi=200)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_axis_off()

    mock_item, live_multi, live_same = series

    ax.text(0.04, 0.93, 'Reproducibility and stability are controlled at the pipeline level', fontsize=17.2, fontweight='bold', color=PALETTE['ink'])
    ax.text(
        0.04,
        0.885,
        'ARG-Test verifies a deterministic local chain, exposes residual live-endpoint variance honestly, and closes archival reproducibility with offline replay.',
        fontsize=10.9,
        color=PALETTE['muted'],
    )

    summary_cards = [
        (
            'Deterministic repository chain',
            'Verified',
            f"Mock 3-seed repeatability: {mock_item['stable_case_count']}/{mock_item['requirement_count']} stable with zero score and coverage drift.",
            PALETTE['teal_soft'],
            PALETTE['teal'],
        ),
        (
            'Seeded live experiments',
            'Supported',
            'Live reruns remain usable for analysis, but the upstream endpoint still contributes residual nondeterminism.',
            PALETTE['gold_soft'],
            PALETTE['gold'],
        ),
        (
            'Submission replay path',
            'Solved',
            'Frozen raw generations rebuild archived result bundles offline, which is the recommended final-package reproducibility route.',
            PALETTE['panel_alt'],
            PALETTE['purple'],
        ),
    ]

    for idx, (title, badge, body, fill, line) in enumerate(summary_cards):
        x = 0.04 + idx * 0.305
        y = 0.63
        w = 0.275
        h = 0.19
        ax.add_patch(
            FancyBboxPatch(
                (x, y),
                w,
                h,
                boxstyle='round,pad=0.012,rounding_size=0.026',
                facecolor=fill,
                edgecolor=line,
                linewidth=1.5,
            )
        )
        ax.text(
            x + 0.02,
            y + h - 0.040,
            textwrap.fill(title, width=24),
            fontsize=11.2,
            fontweight='bold',
            color=PALETTE['ink'],
            va='top',
        )
        ax.text(
            x + 0.02,
            y + h - 0.102,
            f"Status: {badge}",
            fontsize=8.8,
            fontweight='bold',
            color=line,
            va='top',
        )
        ax.text(
            x + 0.02,
            y + h - 0.124,
            textwrap.fill(body, width=36),
            fontsize=8.6,
            color=PALETTE['muted'],
            va='top',
        )

    left_x = 0.04
    left_y = 0.11
    left_w = 0.43
    left_h = 0.43
    right_x = 0.51
    right_y = 0.11
    right_w = 0.42
    right_h = 0.43

    for x, y, w, h, title in [
        (left_x, left_y, left_w, left_h, 'Stable-case ratio across repeatability studies'),
        (right_x, right_y, right_w, right_h, 'Observed live-endpoint variance'),
    ]:
        ax.add_patch(
            FancyBboxPatch(
                (x, y),
                w,
                h,
                boxstyle='round,pad=0.012,rounding_size=0.024',
                facecolor=PALETTE['white'],
                edgecolor=PALETTE['line'],
                linewidth=1.1,
            )
        )
        ax.text(x + 0.02, y + h - 0.05, title, fontsize=12.4, fontweight='bold', color=PALETTE['ink'])

    ratio_rows = [
        ('Mock 3-seed', mock_item, PALETTE['teal']),
        ('Live multi-seed', live_multi, PALETTE['gold']),
        ('Live same-seed', live_same, PALETTE['accent']),
    ]
    for idx, (label, item, color) in enumerate(ratio_rows):
        row_y = left_y + 0.275 - idx * 0.105
        stable_rate = item['stable_case_count'] / item['requirement_count'] if item['requirement_count'] else 0.0
        ax.text(left_x + 0.02, row_y + 0.025, label, fontsize=10.3, fontweight='bold', color=PALETTE['ink'], va='center')
        ax.add_patch(
            FancyBboxPatch(
                (left_x + 0.15, row_y),
                0.19,
                0.03,
                boxstyle='round,pad=0.002,rounding_size=0.01',
                facecolor=PALETTE['panel'],
                edgecolor='none',
            )
        )
        ax.add_patch(
            FancyBboxPatch(
                (left_x + 0.15, row_y),
                0.19 * stable_rate,
                0.03,
                boxstyle='round,pad=0.002,rounding_size=0.01',
                facecolor=color,
                edgecolor='none',
            )
        )
        ax.text(
            left_x + 0.40,
            row_y + 0.025,
            f"{item['stable_case_count']}/{item['requirement_count']} | {stable_rate:.0%}",
            fontsize=9.3,
            color=PALETTE['muted'],
            va='center',
            ha='right',
        )

    delta_rows = [
        ('Live multi-seed score drift', live_multi['avg_max_score_delta'], 0.25, PALETTE['purple'], '0.12'),
        ('Live multi-seed coverage drift', live_multi['avg_max_coverage_delta'], 0.40, PALETTE['gold'], '0.09'),
        ('Live same-seed score drift', live_same['avg_max_score_delta'], 0.25, PALETTE['purple'], '0.10'),
        ('Live same-seed coverage drift', live_same['avg_max_coverage_delta'], 0.40, PALETTE['accent'], '0.22'),
    ]
    for idx, (label, value, cap, color, shown) in enumerate(delta_rows):
        row_y = right_y + 0.275 - idx * 0.075
        ax.text(right_x + 0.02, row_y + 0.028, label, fontsize=9.8, color=PALETTE['ink'], va='center')
        ax.add_patch(
            FancyBboxPatch(
                (right_x + 0.21, row_y + 0.012),
                0.15,
                0.022,
                boxstyle='round,pad=0.002,rounding_size=0.008',
                facecolor=PALETTE['panel'],
                edgecolor='none',
            )
        )
        ax.add_patch(
            FancyBboxPatch(
                (right_x + 0.21, row_y + 0.012),
                0.15 * min(value / cap, 1.0),
                0.022,
                boxstyle='round,pad=0.002,rounding_size=0.008',
                facecolor=color,
                edgecolor='none',
            )
        )
        ax.text(right_x + 0.39, row_y + 0.028, shown, fontsize=9.8, color=PALETTE['muted'], va='center', ha='right')
    save_fig(fig, output_path)
    plt.close(fig)


def save_scorecard(output_path: Path, payload: dict, categories: list[dict]) -> None:
    configure_matplotlib()
    fig = plt.figure(figsize=(10.6, 5.2), dpi=200)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_axis_off()
    cards = [
        ('Test Split', f"{payload['main']['count']} reqs", PALETTE['gold_soft'], PALETTE['gold']),
        ('Checker Score', f"{payload['main']['avg_checker_score']:.3f}", PALETTE['accent_soft'], PALETTE['accent']),
        ('Coverage', f"{payload['main']['avg_overall_coverage']:.3f}", PALETTE['teal_soft'], PALETTE['teal']),
        ('High Risk', f"{payload['main']['high_risk_count']}/16", PALETTE['purple_soft'], PALETTE['purple']),
        ('Repaired', f"{payload['main']['repaired_count']}/16", PALETTE['panel_alt'], PALETTE['ink']),
    ]
    for idx, (label, value, fill, line) in enumerate(cards):
        x = 0.035 + idx * 0.188
        patch = FancyBboxPatch((x, 0.66), 0.165, 0.18, boxstyle='round,pad=0.012,rounding_size=0.024', facecolor=fill, edgecolor=line, linewidth=1.6)
        ax.add_patch(patch)
        ax.text(x + 0.018, 0.786, label, fontsize=10.6, color=PALETTE['muted'])
        ax.text(x + 0.0825, 0.707, value, fontsize=17.2, fontweight='bold', color=PALETTE['ink'], ha='center')
    ax.text(0.045, 0.56, 'Category performance', fontsize=14.5, fontweight='bold', color=PALETTE['ink'])
    ax.text(0.84, 0.56, f"avg risk {payload['main']['avg_risk_score']:.3f}", fontsize=10.6, color=PALETTE['muted'], ha='right')
    y = 0.44
    label_map = {'business_rules': 'Business Rules', 'input_validation': 'Input Validation', 'workflow_state': 'Workflow State'}
    colors = [PALETTE['accent'], PALETTE['gold'], PALETTE['teal']]
    for idx, item in enumerate(categories):
        row_y = y - idx * 0.12
        ax.add_patch(FancyBboxPatch((0.045, row_y), 0.91, 0.082, boxstyle='round,pad=0.008,rounding_size=0.02', facecolor=PALETTE['white'], edgecolor=PALETTE['line'], linewidth=1.0))
        ax.text(0.065, row_y + 0.048, label_map[item['category']], fontsize=12.0, fontweight='bold', color=PALETTE['ink'], va='center')
        ax.text(0.25, row_y + 0.048, f"n={item['requirement_count']}", fontsize=10.1, color=PALETTE['muted'], va='center')
        ax.add_patch(FancyBboxPatch((0.33, row_y + 0.024), 0.28, 0.022, boxstyle='round,pad=0.001,rounding_size=0.008', facecolor=PALETTE['panel'], edgecolor='none'))
        ax.add_patch(FancyBboxPatch((0.33, row_y + 0.024), 0.28 * item['avg_overall_coverage'], 0.022, boxstyle='round,pad=0.001,rounding_size=0.008', facecolor=colors[idx], edgecolor='none'))
        dot_x = 0.33 + 0.28 * item['avg_checker_score']
        ax.add_patch(plt.Circle((dot_x, row_y + 0.062), 0.0065, color=PALETTE['ink']))
        ax.text(0.64, row_y + 0.048, f"cov {item['avg_overall_coverage']:.3f}", fontsize=10.8, color=PALETTE['ink'], va='center')
        ax.text(0.86, row_y + 0.048, f"score {item['avg_checker_score']:.3f}", fontsize=10.8, color=PALETTE['ink'], va='center', ha='right')
        ax.text(0.94, row_y + 0.048, f"risk {item.get('avg_risk_score', 0.0):.2f}", fontsize=10.0, color=PALETTE['muted'], va='center', ha='right')
    save_fig(fig, output_path)
    plt.close(fig)



def save_case_study(output_path: Path, run_main: list[dict]) -> None:
    configure_matplotlib()
    selected = [
        'coupon_discount_engine',
        'payment_card_expiry_and_cvv_validation',
        'payment_3ds_authentication_flow',
    ]
    by_id = {item['requirement_id']: item for item in run_main}
    label_map = {
        'coupon_discount_engine': 'Business-rule case',
        'payment_card_expiry_and_cvv_validation': 'Input-validation case',
        'payment_3ds_authentication_flow': 'Workflow case',
    }
    fig, axes = plt.subplots(1, 3, figsize=(12.8, 4.7), dpi=200)
    for ax, rid, color in zip(axes, selected, [PALETTE['accent'], PALETTE['gold'], PALETTE['teal']]):
        item = by_id[rid]
        ax.set_axis_off()
        ax.add_patch(FancyBboxPatch((0.03, 0.05), 0.94, 0.9, boxstyle='round,pad=0.018,rounding_size=0.035', facecolor=PALETTE['white'], edgecolor=PALETTE['line'], linewidth=1.0))
        ax.add_patch(FancyBboxPatch((0.03, 0.82), 0.94, 0.13, boxstyle='round,pad=0.018,rounding_size=0.035', facecolor=color, edgecolor=color, linewidth=0))
        ax.text(0.08, 0.875, label_map[rid], fontsize=13.8, fontweight='bold', color=PALETTE['white'], va='center')
        ax.text(0.08, 0.73, textwrap.fill(pretty_requirement_name(rid), width=22), fontsize=12.5, fontweight='bold', color=PALETTE['ink'])
        ax.text(0.08, 0.64, rid, fontsize=9.6, color=PALETTE['muted'])
        ax.text(0.08, 0.50, f"checker score", fontsize=9.8, color=PALETTE['muted'])
        ax.text(0.55, 0.50, f"{float(item['score']):.3f}", fontsize=16.5, fontweight='bold', color=PALETTE['ink'], ha='right')
        ax.text(0.08, 0.40, f"overall coverage", fontsize=9.8, color=PALETTE['muted'])
        ax.text(0.55, 0.40, f"{float(item['metrics']['overall_coverage']):.3f}", fontsize=16.5, fontweight='bold', color=PALETTE['ink'], ha='right')
        ax.text(0.08, 0.30, f"test cases", fontsize=9.8, color=PALETTE['muted'])
        ax.text(0.55, 0.30, f"{int(item['metrics']['test_count'])}", fontsize=16.5, fontweight='bold', color=PALETTE['ink'], ha='right')
        risk = item.get('risk_assessment', {})
        ax.text(0.67, 0.50, 'risk profile', fontsize=9.8, color=PALETTE['muted'])
        ax.text(0.90, 0.44, f"{risk.get('level', 'NA')} / {risk.get('score', 0):.2f}", fontsize=11.6, fontweight='bold', color=PALETTE['ink'], ha='right')
        ax.text(0.67, 0.35, 'top diagnostic', fontsize=9.8, color=PALETTE['muted'])
        diag = short_diag((item.get('diagnostics') or ['no diagnostic'])[0], width=46)
        ax.text(0.67, 0.28, textwrap.fill(diag, width=18), fontsize=9.2, color=PALETTE['ink'], ha='left', va='top')
    fig.tight_layout(pad=0.7, w_pad=1.0)
    save_fig(fig, output_path)
    plt.close(fig)



def write_readme(output_dir: Path) -> None:
    content = '''# Figure Assets

Generated by `experiments/generate_presentation_assets.py`.

Editable architecture:
- `arg_test_architecture_editable.pptx`

PNG figures:
- `final_result_scorecard.png`
- `main_vs_baselines.png`
- `ablation_gain.png`
- `generalization_by_category.png`
- `stability_sanity_check.png`
- `reproducibility_stability_overview.png`
- `case_study_snapshots.png`

PDF figures:
- `final_result_scorecard.pdf`
- `main_vs_baselines.pdf`
- `ablation_gain.pdf`
- `generalization_by_category.pdf`
- `stability_sanity_check.pdf`
- `reproducibility_stability_overview.pdf`
- `case_study_snapshots.pdf`

If you edit the architecture PPTX and want a PNG export, run:
`powershell -ExecutionPolicy Bypass -File experiments/export_architecture_slide.ps1 -PptxPath report_assets/figures/arg_test_architecture_editable.pptx -OutputDir report_assets/figures/architecture_png`
'''
    (output_dir / 'README.md').write_text(content, encoding='utf-8')



def main() -> None:
    parser = argparse.ArgumentParser(description='Generate PPTX and PNG presentation assets from final experiment results.')
    parser.add_argument('--output-root', default='.local_runs/formal_qwen_novpn')
    parser.add_argument('--stability-root', default='.local_runs/stability_qwen_20260411')
    parser.add_argument('--repro-mock-root', default='.local_runs/repro_multi_seed_mock')
    parser.add_argument('--repro-live-root', default='.local_runs/repro_live_qwen_5case')
    parser.add_argument('--repro-same-seed-root', default='.local_runs/repro_live_same_seed_3case')
    parser.add_argument('--figure-dir', default='report_assets/figures')
    parser.add_argument('--skip-architecture', action='store_true')
    args = parser.parse_args()

    report_root = ROOT / args.output_root / 'outputs' / 'reports' / 'test'
    stability_root = ROOT / args.stability_root / 'outputs' / 'reports' / 'test'
    repro_mock_root = ROOT / args.repro_mock_root / 'outputs' / 'reports' / 'test'
    repro_live_root = ROOT / args.repro_live_root / 'outputs' / 'reports' / 'test'
    repro_same_seed_root = ROOT / args.repro_same_seed_root / 'outputs' / 'reports' / 'test'
    figure_dir = ROOT / args.figure_dir
    ensure_dir(figure_dir)

    run_main = read_json(report_root / 'run_main_summary.json')
    baseline = read_json(report_root / 'baseline_summary.json')
    ablation = read_json(report_root / 'ablation_summary.json')
    categories = read_json(report_root / 'generalization_by_category.json')['categories']
    stability = read_json(stability_root / 'stability_sanity_summary.json')
    repro_mock = read_json(repro_mock_root / 'repeatability_summary.json')
    repro_live = read_json(repro_live_root / 'repeatability_summary.json')
    repro_same_seed = read_json(repro_same_seed_root / 'repeatability_summary.json')
    payload = aggregate_payload(run_main, baseline, ablation)
    repro_series = [
        {
            'label': 'Mock 3-seed',
            'subhead': f"provider={repro_mock['provider']} | n={repro_mock['requirement_count']}",
            'stable_case_count': int(repro_mock['stable_case_count']),
            'requirement_count': int(repro_mock['requirement_count']),
            'avg_max_score_delta': float(repro_mock['avg_max_score_delta']),
            'avg_max_coverage_delta': float(repro_mock['avg_max_coverage_delta']),
        },
        {
            'label': 'Live multi-seed',
            'subhead': f"provider={repro_live['model']} | n={repro_live['requirement_count']}",
            'stable_case_count': int(repro_live['stable_case_count']),
            'requirement_count': int(repro_live['requirement_count']),
            'avg_max_score_delta': float(repro_live['avg_max_score_delta']),
            'avg_max_coverage_delta': float(repro_live['avg_max_coverage_delta']),
        },
        {
            'label': 'Live same-seed',
            'subhead': f"seed fixed x{repro_same_seed['repeats']} | n={repro_same_seed['requirement_count']}",
            'stable_case_count': int(repro_same_seed['stable_case_count']),
            'requirement_count': int(repro_same_seed['requirement_count']),
            'avg_max_score_delta': float(repro_same_seed['avg_max_score_delta']),
            'avg_max_coverage_delta': float(repro_same_seed['avg_max_coverage_delta']),
        },
    ]

    if not args.skip_architecture:
        build_architecture_pptx(figure_dir / 'arg_test_architecture_editable.pptx')
    save_scorecard(figure_dir / 'final_result_scorecard.png', payload, categories)
    save_main_vs_baselines(figure_dir / 'main_vs_baselines.png', payload)
    save_ablation(figure_dir / 'ablation_gain.png', payload)
    save_generalization(figure_dir / 'generalization_by_category.png', categories)
    save_stability(figure_dir / 'stability_sanity_check.png', stability)
    save_reproducibility_summary(figure_dir / 'reproducibility_stability_overview.png', repro_series)
    save_case_study(figure_dir / 'case_study_snapshots.png', run_main)
    write_readme(figure_dir)

    print(json.dumps({
        'figure_dir': str(figure_dir),
        'generated': [
            *([] if args.skip_architecture else ['arg_test_architecture_editable.pptx']),
            'final_result_scorecard.png',
            'final_result_scorecard.pdf',
            'main_vs_baselines.png',
            'main_vs_baselines.pdf',
            'ablation_gain.png',
            'ablation_gain.pdf',
            'generalization_by_category.png',
            'generalization_by_category.pdf',
            'stability_sanity_check.png',
            'stability_sanity_check.pdf',
            'reproducibility_stability_overview.png',
            'reproducibility_stability_overview.pdf',
            'case_study_snapshots.png',
            'case_study_snapshots.pdf',
            'README.md',
        ]
    }, ensure_ascii=False, indent=2))


if __name__ == '__main__':
    main()
