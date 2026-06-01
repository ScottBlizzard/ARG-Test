from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except Exception:  # pragma: no cover - optional fallback
    Image = None


ROOT = Path(__file__).resolve().parents[2]
ASSET_ROOT = ROOT / "07_PPT_Assets_For_Luowu"
OUT_PATH = ASSET_ROOT / "Group7_Final_PPT_MiniShop_Rebuilt_Editable.pptx"

SLIDE_W = 13.333333
SLIDE_H = 7.5
PX_PER_IN = 96.0


THEME = {
    "bg": "FFFDFC",
    "paper": "FFF7F6",
    "card": "FFFFFF",
    "blush": "F6DAD7",
    "blush_dark": "E7B9B4",
    "rose": "D88F88",
    "ink": "253042",
    "body": "5C6470",
    "muted": "8A9098",
    "green": "2E8B73",
    "gold": "D2A63E",
    "blue": "7FA9C9",
    "line": "E8CECA",
}


MEMBERS = [
    "Yi Xu 2351441",
    "Xiang Wang 2351039",
    "Fengxuan Kang 2350283",
    "Luowu Zhang 2352746",
    "Yiwei Chen 2350217",
]


def px(value: float) -> int:
    return int(Inches(value / PX_PER_IN))


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.replace("#", "")
    return RGBColor.from_string(hex_color)


def set_fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill, line, line_width=1.0, radius=False):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        px(x),
        px(y),
        px(w),
        px(h),
    )
    set_fill(shp, fill)
    set_line(shp, line, line_width)
    if radius:
        shp.adjustments[0] = 0.08
    return shp


def add_line(slide, x, y, w, h=1.5, color=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px(x), px(y), px(w), px(h))
    fill = color or THEME["line"]
    set_fill(shp, fill)
    shp.line.fill.background()
    return shp


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=20,
    color=None,
    bold=False,
    font="Aptos",
    align="left",
    valign="top",
):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)
    lines = str(text).split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color or THEME["body"])
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(align, PP_ALIGN.LEFT)
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    return tb


def add_bullets(slide, items, x, y, w, *, size=18, line_height=28, text_color=None, bullet_color=None):
    for idx, item in enumerate(items):
        yy = y + idx * line_height
        add_text(
            slide,
            "•",
            x,
            yy,
            16,
            18,
            size=size,
            color=bullet_color or THEME["rose"],
            bold=True,
        )
        add_text(
            slide,
            item,
            x + 18,
            yy,
            w - 18,
            line_height + 10,
            size=size,
            color=text_color or THEME["body"],
        )


def add_chip(slide, text, x, y, w=190):
    add_rect(slide, x, y, w, 28, THEME["paper"], THEME["line"], 1.0, radius=True)
    add_text(slide, text, x + 10, y + 6, w - 20, 16, size=11, color=THEME["ink"], bold=True)


def add_bg(slide):
    add_rect(slide, 0, 0, 1280, 720, THEME["bg"], THEME["bg"], 0)
    add_rect(slide, 0, 0, 28, 720, THEME["blush"], THEME["blush"], 0)
    add_rect(slide, 0, 702, 1280, 18, "F6E4E2", "F6E4E2", 0)


def add_title(slide, title_text, kicker=None):
    if kicker:
        add_text(slide, kicker, 60, 34, 260, 20, size=12, color=THEME["rose"], bold=True)
        top = 52
    else:
        top = 42
    add_text(slide, title_text, 60, top, 1120, 62, size=31, color=THEME["rose"], font="Georgia")


def add_picture_contain(slide, image_path: Path, x, y, w, h):
    if not image_path.exists():
        return None
    box_w = px(w)
    box_h = px(h)
    left = px(x)
    top = px(y)
    if Image and image_path.suffix.lower() not in {".svg"}:
        with Image.open(image_path) as img:
            img_w, img_h = img.size
        scale = min(box_w / img_w, box_h / img_h)
        draw_w = int(img_w * scale)
        draw_h = int(img_h * scale)
        draw_x = left + int((box_w - draw_w) / 2)
        draw_y = top + int((box_h - draw_h) / 2)
        return slide.shapes.add_picture(str(image_path), draw_x, draw_y, width=draw_w, height=draw_h)
    return slide.shapes.add_picture(str(image_path), left, top, width=box_w, height=box_h)


def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, "Software Testing 42036101 Final Project", 360, 212, 560, 28, size=16, color=THEME["muted"], align="center")
    add_text(slide, "ARG-Test", 220, 110, 840, 70, size=42, color=THEME["ink"], bold=True, font="Georgia", align="center")
    add_text(
        slide,
        "Auditable and Risk-Aware\nRequirement-Driven Black-Box Test Generation",
        180,
        154,
        920,
        70,
        size=22,
        color=THEME["body"],
        align="center",
    )
    add_rect(slide, 420, 260, 440, 34, "FFF9F8", THEME["line"], 1.2, radius=True)
    add_text(slide, "Team ID: Group 7", 420, 268, 440, 18, size=15, color=THEME["ink"], bold=True, align="center")
    add_rect(slide, 205, 355, 870, 210, "FFFFFF", THEME["line"], 1.0, radius=True)
    add_text(slide, "Members", 255, 385, 160, 22, size=16, color=THEME["rose"], bold=True)
    positions = [(255, 430), (255, 464), (255, 498), (615, 430), (615, 464)]
    for member, (x, y) in zip(MEMBERS, positions):
        add_text(slide, member, x, y, 300, 22, size=18, color=THEME["body"])
    add_text(slide, "Tool: ARG-Test | Target application: MiniShop Checkout", 205, 592, 870, 20, size=14, color=THEME["muted"], align="center")


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Assignment Goal and Final Object Separation")
    add_rect(slide, 82, 160, 1116, 110, "FFFFFF", THEME["line"], 1.0, radius=True)
    add_text(
        slide,
        "Assignment goal: build an AI-driven AutoTestDesign tool and use it to test an independent target application.",
        110,
        195,
        1060,
        50,
        size=24,
        color=THEME["ink"],
        align="center",
    )
    boxes = [
        (96, "TOOL", "ARG-Test", "Requirement analysis\nRisk scoring\nSystematic test design", THEME["blush"], 22, 16, 52, 70),
        (454, "TARGET APPLICATION", "MiniShop\nCheckout", "Promotion\nShipping / Tax\nPayment / Pickup\nCheckout orchestration", "EAF5F1", 21, 13, 60, 74),
        (812, "SELECTED MODULE", "coupon_discount_\nengine", "Promotion service coupon engine\nDetailed executable evidence", "FFF5E8", 18, 13, 64, 58),
    ]
    for x, label, title_text, note, fill, title_size, note_size, note_top, note_height in boxes:
        add_rect(slide, x, 320, 292, 220, fill, THEME["line"], 1.4, radius=True)
        add_chip(slide, label, x + 16, 338, 120 if label == "TOOL" else (210 if label == "TARGET APPLICATION" else 150))
        add_text(slide, title_text, x + 18, 390, 256, 62, size=title_size, color=THEME["ink"], bold=True, font="Georgia")
        add_text(slide, note, x + 18, 390 + note_top, 250, note_height, size=note_size, color=THEME["body"])
    add_text(slide, "This separation is the key correction from the earlier version.", 82, 590, 1116, 24, size=17, color=THEME["rose"], bold=True, align="center")


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "MiniShop Checkout: Target Application Scope")
    add_rect(slide, 84, 142, 430, 430, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_text(
        slide,
        "MiniShop Checkout is a compact\ne-commerce checkout prototype\nused as the independent target\napplication in this project.",
        114,
        185,
        370,
        108,
        size=18,
        color=THEME["ink"],
        font="Georgia",
    )
    add_line(slide, 114, 300, 360)
    add_text(slide, "Implemented scope", 114, 320, 180, 24, size=16, color=THEME["rose"], bold=True)
    add_bullets(
        slide,
        [
            "Promotion and coupon handling",
            "Shipping-fee calculation",
            "Tax and order-total calculation",
            "Payment-card validation",
            "Pickup-contact validation",
            "Checkout preview orchestration",
        ],
        114,
        358,
        360,
        size=15,
        line_height=35,
    )
    add_rect(slide, 550, 142, 648, 430, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_text(slide, "Component map", 580, 175, 200, 24, size=16, color=THEME["rose"], bold=True)
    comps = [
        ("Promotion Service", "coupon rules and restrictions"),
        ("Shipping Service", "fee tiers and shipping thresholds"),
        ("Tax Service", "taxable amount and order total"),
        ("Payment Validation", "card number, expiry, CVV, brand rules"),
        ("Pickup Validation", "station ID, phone, pickup code"),
        ("Checkout Service", "end-to-end preview orchestration"),
    ]
    for index, (head, note) in enumerate(comps):
        row = index // 2
        col = index % 2
        x = 580 + col * 300
        y = 218 + row * 110
        fill = "FFF7F6" if index % 2 == 0 else "F7FBFA"
        add_rect(slide, x, y, 260, 88, fill, THEME["line"], 1.0, radius=True)
        add_text(slide, head, x + 16, y + 14, 228, 24, size=16, color=THEME["ink"], bold=True)
        add_text(slide, note, x + 16, y + 42, 228, 28, size=13, color=THEME["body"])
    add_text(
        slide,
        "Out of scope: refund workflows, inventory sync, external payment gateway integration",
        84,
        610,
        1116,
        20,
        size=16,
        color=THEME["muted"],
        align="center",
    )


def slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Why Plain LLM Is Not Enough")
    add_text(slide, "Plain LLM outputs can be fluent, but they are not reliably auditable.", 82, 128, 1110, 24, size=20, color=THEME["body"])
    add_rect(slide, 88, 176, 500, 360, "FFF9F8", THEME["line"], 1.3, radius=True)
    add_rect(slide, 692, 176, 500, 360, "F9FCFB", THEME["line"], 1.3, radius=True)
    add_chip(slide, "PLAIN LLM OUTPUT", 108, 196, 150)
    add_chip(slide, "ARG-TEST STRUCTURED TRACE", 712, 196, 210)
    add_text(
        slide,
        "Try SAVE10 on a valid order.\nTry an expired coupon.\nCheck that the discount is correct.\n\nLooks fluent, but gives no clear\nboundary, conflict, or traceability plan.",
        116,
        244,
        438,
        220,
        size=22,
        color=THEME["body"],
    )
    add_text(
        slide,
        "Analysis: thresholds 30 / 50 / 100,\npremium-only SAVE20, expiry, sale-item conflict\n\nPattern: EP + BVA + Decision Table\n\nVerification: boundary, invalid, and exclusivity\nobligations are explicitly checked\n\nFinal tests: subtotal=50, expired coupon,\nSAVE20 + sale items, coupon stacking rejected",
        720,
        244,
        438,
        230,
        size=18,
        color=THEME["body"],
    )
    add_text(slide, "The problem is not only correctness. The deeper issue is auditability.", 120, 585, 1040, 30, size=22, color=THEME["rose"], bold=True, align="center")


def slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "ARG-Test Final System Overview")
    add_picture_contain(slide, ASSET_ROOT / "figures_png" / "arg_test_architecture_final.png", 84, 132, 760, 500)
    add_rect(slide, 882, 150, 280, 444, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_chip(slide, "PIPELINE", 902, 172, 86)
    add_bullets(
        slide,
        [
            "ARG-Test is a pipeline, not a single prompt",
            "Structured generation",
            "Parser and schema gate",
            "Technique-aware checker",
            "Rerank and targeted repair",
            "Export, evaluation, and reproducible reporting",
        ],
        904,
        216,
        232,
        size=18,
        line_height=48,
    )


def slide_06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Structured Trace and Contract Checking")
    labels = ["Analysis", "Pattern", "Steps", "Verification", "Final Answer"]
    xs = [96, 306, 516, 726, 936]
    widths = [170, 170, 170, 170, 180]
    fills = ["FFF8F7", "FFF8F7", "FFF8F7", "FFF8F7", "F2D6D3"]
    for idx, label in enumerate(labels):
        add_rect(slide, xs[idx], 172, widths[idx], 96, fills[idx], THEME["line"], 1.2, radius=True)
        size = 22 if idx < 3 else 18
        add_text(slide, label, xs[idx] + 10, 202, widths[idx] - 20, 36, size=size, color=THEME["ink"], bold=True, align="center", font="Georgia")
        if idx < len(labels) - 1:
            add_line(slide, xs[idx] + widths[idx], 219, xs[idx + 1] - (xs[idx] + widths[idx]))
            add_text(slide, ">", xs[idx] + widths[idx] + 14, 205, 18, 20, size=20, color=THEME["rose"], bold=True)
    add_rect(slide, 132, 328, 1016, 114, "FFF4F3", THEME["line"], 1.2, radius=True)
    add_text(slide, "Technique-Aware Contract Checking", 160, 350, 960, 22, size=18, color=THEME["rose"], bold=True, align="center")
    chips = ["EP Checker", "BVA Checker", "Decision Checker", "State Checker"]
    chip_x = [184, 418, 652, 896]
    chip_w = [150, 150, 180, 150]
    for cx, cw, chip in zip(chip_x, chip_w, chips):
        add_rect(slide, cx, 382, cw, 34, "FFFFFF", THEME["line"], 1.0, radius=True)
        add_text(slide, chip, cx + 8, 390, cw - 16, 18, size=14, color=THEME["ink"], bold=True, align="center")
    add_rect(slide, 96, 478, 1088, 126, "FFFFFF", THEME["line"], 1.1, radius=True)
    add_text(slide, "Typed trace: Analysis -> Pattern -> Steps -> Verification -> FinalAnswer", 128, 508, 980, 26, size=22, color=THEME["ink"], bold=True)
    add_text(slide, "The model output becomes a typed artifact that can be parsed, checked, and repaired.", 128, 550, 980, 24, size=18, color=THEME["body"])


def slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Requirement Closure and Interactive Review")
    add_rect(slide, 86, 144, 510, 462, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_text(slide, "Mandatory requirement closure", 114, 176, 300, 22, size=16, color=THEME["rose"], bold=True)
    add_text(slide, "FR 1.0, FR 1.1, FR 2.0, FR 3.0, FR 6.0", 114, 208, 430, 30, size=24, color=THEME["ink"], bold=True)
    add_text(slide, "Extra-credit closure", 114, 258, 250, 22, size=16, color=THEME["rose"], bold=True)
    add_text(slide, "FR 4.0, FR 5.0, FR 7.0", 114, 290, 430, 30, size=24, color=THEME["ink"], bold=True)
    add_text(slide, "Interactive review surface", 114, 342, 260, 22, size=16, color=THEME["rose"], bold=True)
    add_bullets(
        slide,
        [
            "Direct Input",
            "CSV Batch",
            "State Model",
            "Formal Evidence",
            "Inspect outputs, revise guidance, rerun the pipeline",
        ],
        114,
        378,
        430,
        size=18,
        line_height=36,
    )
    add_picture_contain(slide, ASSET_ROOT / "frontend_screenshots" / "web_demo_direct_frozen_replay.png", 634, 164, 500, 344)
    add_text(slide, "The tester remains in the loop through inspection, revision, and rerun.", 640, 542, 490, 26, size=18, color=THEME["body"], align="center")


def slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Risk Analysis for MiniShop Checkout")
    add_rect(slide, 88, 150, 350, 426, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_text(slide, "Risk scoring method", 116, 178, 240, 22, size=16, color=THEME["rose"], bold=True)
    add_text(slide, "Risk Priority =\nImpact x Likelihood x Detectability", 116, 220, 290, 86, size=26, color=THEME["ink"], bold=True, font="Georgia")
    add_line(slide, 116, 328, 290)
    add_text(slide, "Priority bands", 116, 350, 180, 20, size=16, color=THEME["rose"], bold=True)
    add_bullets(slide, ["High >= 60", "Medium = 36 to 59", "Low <= 35"], 116, 382, 260, size=19, line_height=42)
    add_rect(slide, 478, 150, 710, 426, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_text(slide, "Highest-risk areas", 510, 178, 220, 22, size=16, color=THEME["rose"], bold=True)
    risk_rows = [
        ("coupon and promotion logic", THEME["blush"]),
        ("shipping and tax calculation", "FDEFCF"),
        ("payment-card validation", "E2F2ED"),
        ("checkout orchestration", "EDF3FB"),
        ("pickup validation", "F8F1FF"),
    ]
    for idx, (label, fill) in enumerate(risk_rows):
        y = 216 + idx * 62
        add_rect(slide, 510, y, 640, 46, fill, THEME["line"], 1.0, radius=True)
        add_text(slide, label, 528, y + 12, 606, 20, size=19, color=THEME["ink"], bold=idx < 4)
    add_text(slide, "High risk drives the most detailed test effort and the selected executable module.", 88, 610, 1100, 22, size=17, color=THEME["body"], align="center")


def slide_09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Test Plan: Scope, Architecture, and Test Suites")
    starts = [88, 477, 866]
    headings = ["Scope", "Architecture", "Planned Suites"]
    fills = ["FFF7F6", "F9FCFB", "FFFBF4"]
    bodies = [
        "Promotion and pricing\nShipping and tax\nPayment validation\nPickup validation\nCheckout orchestration",
        "Checkout Service\nPromotion Service\nShipping Service\nTax Service\nPayment Validation\nPickup Validation",
        "Promotion suite\nShipping and tax suite\nPayment validation suite\nPickup validation suite\nCheckout orchestration suite\nDetailed module suite",
    ]
    for x, heading, fill, body in zip(starts, headings, fills, bodies):
        add_rect(slide, x, 164, 326, 404, fill, THEME["line"], 1.2, radius=True)
        add_text(slide, heading, x + 22, 194, 200, 24, size=18, color=THEME["rose"], bold=True)
        add_text(slide, body, x + 22, 242, 282, 290, size=22, color=THEME["ink"])
    add_text(slide, "The test plan covers the target application, not the AutoTestDesign tool itself.", 88, 610, 1100, 20, size=17, color=THEME["body"], align="center")


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Test Plan: Schedule, Team, Framework, and Cost")
    quads = [
        (88, 166, "Schedule", "FFF7F6", "freeze target-application scope\ncomplete risk analysis\nreview generated suites\nexecute detailed module tests\npackage final evidence", 15),
        (652, 166, "Team", "F9FCFB", "integration\ndocument and PPT\nrequirement assets\nmodule execution support\nevaluation and reproducibility", 15),
        (88, 390, "Framework", "FFFBF4", "pytest + coverage.py\nPython-native stack\nclear black-box and white-box assertions\nsimple coverage integration", 16),
        (652, 390, "Cost", "F8F3FF", "With ARG-Test:\n4.5 to 7.0 person-days\n\nManual baseline:\n7.5 to 10.0 person-days", 16),
    ]
    for x, y, label, fill, body, body_size in quads:
        add_rect(slide, x, y, 460, 180, fill, THEME["line"], 1.2, radius=True)
        add_text(slide, label, x + 20, y + 18, 180, 22, size=18, color=THEME["rose"], bold=True)
        add_text(slide, body, x + 20, y + 52, 420, 110, size=body_size, color=THEME["ink"])
    add_text(slide, "This slide closes the teacher-required plan items: schedule, team, framework, and cost.", 88, 636, 1100, 20, size=16, color=THEME["muted"], align="center")


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Experimental Setup and Evaluation Protocol")
    add_text(slide, "Tool-level evaluation on the broader requirement corpus", 88, 116, 1100, 20, size=16, color=THEME["body"])
    add_picture_contain(slide, ASSET_ROOT / "slide_experiment_assets" / "slide_experimental_setup_protocol.png", 72, 146, 1136, 520)


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Baselines and Fair Comparison")
    add_picture_contain(slide, ASSET_ROOT / "slide_experiment_assets" / "slide_baselines_fair_comparison.png", 72, 140, 1136, 520)
    add_text(slide, "The rule-based baseline is our own deterministic heuristic non-AI baseline.", 88, 612, 1100, 18, size=16, color=THEME["body"], align="center")


def slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Main Result Scorecard")
    add_picture_contain(slide, ASSET_ROOT / "figures_png" / "final_result_scorecard.png", 84, 146, 790, 500)
    add_rect(slide, 914, 206, 232, 286, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_text(slide, "Frozen test split summary", 934, 236, 190, 22, size=16, color=THEME["rose"], bold=True)
    add_bullets(
        slide,
        [
            "16 held-out requirements",
            "Avg checker score: 0.959",
            "Avg overall coverage: 0.615",
            "Avg tests: 7.312",
            "0 duplicate cases",
            "9 / 16 high-risk requirements prioritized",
        ],
        934,
        276,
        180,
        size=16,
        line_height=34,
    )


def slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Main Comparison Against Baselines")
    add_picture_contain(slide, ASSET_ROOT / "figures_png" / "main_vs_baselines.png", 82, 150, 760, 500)
    add_rect(slide, 876, 188, 282, 378, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_text(slide, "Full ARG-Test outperforms all three baselines on the frozen test split.", 900, 224, 234, 84, size=22, color=THEME["ink"], font="Georgia", bold=True)
    add_bullets(
        slide,
        [
            "Against rule-based: better requirement understanding",
            "Against plain LLM: prompting alone is not enough",
            "Against structured no-checker: checker-guided control adds real value",
        ],
        900,
        334,
        236,
        size=17,
        line_height=58,
    )


def slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Generalization and Ablation")
    add_picture_contain(slide, ASSET_ROOT / "figures_png" / "generalization_by_category.png", 78, 154, 520, 350)
    add_picture_contain(slide, ASSET_ROOT / "figures_png" / "ablation_gain.png", 654, 154, 530, 350)
    add_rect(slide, 78, 534, 1106, 84, "FFFFFF", THEME["line"], 1.1, radius=True)
    add_text(slide, "The method generalizes across business rules, input validation, and workflow-state requirements.", 108, 554, 1000, 24, size=19, color=THEME["ink"])
    add_text(slide, "Checker-guided control greatly improves checker alignment while keeping coverage and test count comparable.", 108, 582, 1000, 22, size=18, color=THEME["body"])


def slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Representative MiniShop Cases")
    starts = [74, 452, 830]
    widths = [338, 338, 378]
    titles = ["coupon_discount_\nengine", "pickup_station_contact_\nvalidation", "payment_card_expiry_and_\ncvv_validation"]
    for i in range(3):
        add_rect(slide, starts[i], 136, widths[i], 500, "FFFFFF", THEME["line"], 1.1, radius=True)
        title_size = 16 if i < 2 else 15
        add_text(slide, titles[i], starts[i] + 16, 160, widths[i] - 32, 54, size=title_size, color=THEME["ink"], bold=True, font="Georgia")
    add_picture_contain(slide, ASSET_ROOT / "slide12_case_images" / "01_coupon_discount_engine_panel.png", 88, 224, 308, 378)
    add_picture_contain(slide, ASSET_ROOT / "slide12_case_images" / "02_pickup_station_contact_validation_panel.png", 466, 224, 308, 378)
    add_chip(slide, "PAYMENT VALIDATION", 850, 224, 154)
    add_text(slide, "Core rules", 854, 268, 120, 20, size=15, color=THEME["rose"], bold=True)
    add_bullets(
        slide,
        [
            "card_number: 13 to 19 digits",
            "expiry_month: 1 to 12",
            "expiry_year: current year to current year + 15",
            "Visa/Mastercard: 3-digit CVV",
            "Amex: 4-digit CVV",
            "Masked numbers are rejected",
        ],
        854,
        298,
        318,
        size=14,
        line_height=36,
    )
    add_line(slide, 854, 530, 320)
    add_text(slide, "Techniques", 854, 546, 120, 18, size=15, color=THEME["rose"], bold=True)
    add_text(slide, "EP + BVA", 854, 574, 200, 24, size=20, color=THEME["ink"], bold=True)
    add_text(
        slide,
        "These cases show that the target application covers business rules and input validation with concrete module boundaries.",
        86,
        650,
        1100,
        18,
        size=16,
        color=THEME["muted"],
        align="center",
    )


def slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Detailed Executable Evidence")
    add_picture_contain(slide, ASSET_ROOT / "frontend_screenshots" / "coupon_module_evidence_scorecard.png", 88, 160, 760, 430)
    add_rect(slide, 882, 174, 284, 404, "FFFFFF", THEME["line"], 1.2, radius=True)
    add_text(slide, "Selected major module of MiniShop Checkout", 906, 204, 236, 42, size=16, color=THEME["rose"], bold=True)
    add_text(slide, "coupon_discount_\nengine", 906, 256, 236, 52, size=18, color=THEME["ink"], font="Georgia", bold=True)
    add_bullets(
        slide,
        [
            "15 module tests passed",
            "100% statement coverage",
            "100% branch coverage",
            "4 / 4 mutants killed",
            "Black-box design + white-box execution + mutation usefulness",
        ],
        906,
        328,
        236,
        size=14,
        line_height=42,
    )
    add_text(slide, "This gives us executable evidence, not only design-level metrics.", 88, 620, 1100, 22, size=18, color=THEME["body"], align="center")


def slide_18(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Reproducibility, NFR, Limitations, and Conclusion")
    add_picture_contain(slide, ASSET_ROOT / "figures_png" / "reproducibility_stability_overview.png", 76, 142, 610, 350)
    add_rect(slide, 720, 152, 456, 186, "FFFFFF", THEME["line"], 1.1, radius=True)
    add_text(slide, "NFR snapshot", 744, 174, 160, 22, size=16, color=THEME["rose"], bold=True)
    add_bullets(
        slide,
        [
            "Usability: Direct Input, CSV Batch, State Model, Formal Evidence",
            "Security: no secret leak found in exported artifacts",
            "Maintainability: modular codebase and rerunnable tests",
        ],
        744,
        206,
        396,
        size=13,
        line_height=42,
    )
    add_rect(slide, 720, 350, 456, 150, "FFF9F8", THEME["line"], 1.1, radius=True)
    add_text(slide, "Limitations", 744, 376, 140, 22, size=16, color=THEME["rose"], bold=True)
    add_bullets(
        slide,
        [
            "MiniShop Checkout is a compact course-project prototype",
            "Only coupon_discount_engine has full white-box execution evidence",
            "Live providers still show residual nondeterminism",
        ],
        744,
        406,
        394,
        size=12,
        line_height=34,
    )
    add_rect(slide, 76, 520, 1100, 126, "FFFFFF", THEME["line"], 1.1, radius=True)
    add_text(
        slide,
        "ARG-Test enables auditable, checked, and reproducible requirement-driven testing for a concrete target application.",
        104,
        546,
        1040,
        46,
        size=20,
        color=THEME["ink"],
        font="Georgia",
        bold=True,
        align="center",
    )
    add_text(
        slide,
        "Tool: ARG-Test | Target application: MiniShop Checkout | Selected module: coupon_discount_engine",
        104,
        604,
        1040,
        18,
        size=15,
        color=THEME["muted"],
        align="center",
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide_fns = [
        slide_01,
        slide_02,
        slide_03,
        slide_04,
        slide_05,
        slide_06,
        slide_07,
        slide_08,
        slide_09,
        slide_10,
        slide_11,
        slide_12,
        slide_13,
        slide_14,
        slide_15,
        slide_16,
        slide_17,
        slide_18,
    ]
    for fn in slide_fns:
        fn(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved editable PPTX to: {OUT_PATH}")


if __name__ == "__main__":
    main()
